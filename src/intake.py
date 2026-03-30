"""Input scanning and content extraction for the essay writer."""

from __future__ import annotations

import base64
import mimetypes
from pathlib import Path

import pymupdf
from docx import Document as DocxDocument
from pptx import Presentation

from src.tools.docx_reader import extract_docx_text

# ---------------------------------------------------------------------------
# File type classification
# ---------------------------------------------------------------------------

_TEXT_EXTENSIONS = {".md", ".txt", ".text", ".rst", ".csv", ".tsv", ".log"}

_PDF_EXTENSIONS = {".pdf"}

_DOCX_EXTENSIONS = {".docx"}

_PPTX_EXTENSIONS = {".pptx"}

_IMAGE_EXTENSIONS = {
    ".png",
    ".jpg",
    ".jpeg",
    ".gif",
    ".bmp",
    ".tiff",
    ".tif",
    ".webp",
    ".svg",
}

_UNSUPPORTED_KNOWN = {
    ".doc": "Old Word binary format — save as .docx first",
    ".rtf": "RTF format — save as .docx or .txt first",
    ".odt": "OpenDocument format — save as .docx first",
    ".ppt": "Old PowerPoint binary format — save as .pptx or .pdf first",
    ".xls": "Excel — save as .csv or .pdf first",
    ".xlsx": "Excel — save as .csv or .pdf first",
    ".pages": "Apple Pages — export as .pdf or .docx first",
    ".key": "Apple Keynote — export as .pdf first",
    ".numbers": "Apple Numbers — export as .csv first",
}

SUPPORTED_EXTENSIONS = (
    _TEXT_EXTENSIONS
    | _PDF_EXTENSIONS
    | _DOCX_EXTENSIONS
    | _PPTX_EXTENSIONS
    | _IMAGE_EXTENSIONS
)


_EXTENSION_CATEGORIES: dict[str, str] = {
    ext: cat
    for cat, exts in [
        ("text", _TEXT_EXTENSIONS),
        ("pdf", _PDF_EXTENSIONS),
        ("docx", _DOCX_EXTENSIONS),
        ("pptx", _PPTX_EXTENSIONS),
        ("image", _IMAGE_EXTENSIONS),
    ]
    for ext in exts
}


def _classify(path: Path) -> str:
    """Return the category of a file: 'text', 'pdf', 'docx', 'pptx', 'image', or 'unsupported'."""
    return _EXTENSION_CATEGORIES.get(path.suffix.lower(), "unsupported")


# ---------------------------------------------------------------------------
# Content extraction
# ---------------------------------------------------------------------------


def _extract_text(path: Path) -> str:
    return path.read_text(encoding="utf-8")


# Minimum average characters per page to consider text extraction usable.
# Below this threshold the PDF is likely scanned/image-based.
_MIN_CHARS_PER_PAGE = 50


def _extract_pdf(path: Path) -> tuple[str | None, list[dict] | None]:
    """Extract content from a PDF, falling back to page images for scanned PDFs.

    Returns (text, None) for text-native PDFs, or (None, image_blocks) for
    scanned/image-based PDFs where text extraction yields too little content.
    """
    doc = pymupdf.open(str(path))
    total = len(doc)

    # First pass: try text extraction
    parts: list[str] = []
    total_chars = 0
    for i, page in enumerate(doc):
        page_text = page.get_text()
        total_chars += len(page_text.strip())
        parts.append(f"--- Page {i + 1} of {total} ---\n{page_text}")

    avg_chars = total_chars / max(total, 1)

    if avg_chars >= _MIN_CHARS_PER_PAGE:
        doc.close()
        return "\n\n".join(parts), None

    # Fallback: render pages as images for multimodal processing
    image_blocks = _render_pdf_pages(doc)
    doc.close()
    return None, image_blocks


def _make_image_block(image_bytes: bytes, mime: str = "image/png") -> dict:
    """Encode raw bytes as a base64 multimodal content block."""
    data = base64.standard_b64encode(image_bytes).decode("ascii")
    return {
        "type": "image_url",
        "image_url": {"url": f"data:{mime};base64,{data}"},
    }


def _render_pdf_pages(doc: pymupdf.Document) -> list[dict]:
    """Render PDF pages as PNG images for multimodal LLM consumption."""
    blocks: list[dict] = []
    for page in doc:
        # Render at 150 DPI — good balance between quality and size
        pix = page.get_pixmap(dpi=150)
        blocks.append(_make_image_block(pix.tobytes("png")))
    return blocks


def _extract_docx(path: Path) -> str:
    return extract_docx_text(DocxDocument(str(path)))


def _extract_pptx(path: Path) -> str:
    """Extract text from a .pptx file — slide titles, body text, and notes."""
    prs = Presentation(str(path))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        slide_parts: list[str] = [f"--- Slide {i} ---"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_parts.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    slide_parts.append(" | ".join(cells))
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                slide_parts.append(f"[Speaker notes: {notes_text}]")
        if len(slide_parts) > 1:  # more than just the header
            parts.append("\n".join(slide_parts))
    return "\n\n".join(parts)


def _image_mime(path: Path) -> str:
    mime, _ = mimetypes.guess_type(str(path))
    return mime or "image/png"


def _load_image_block(path: Path) -> dict:
    """Load an image as a base64 content block for multimodal messages."""
    return _make_image_block(path.read_bytes(), _image_mime(path))


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------


class InputFile:
    """Represents a discovered input file with its extracted content."""

    __slots__ = ("path", "category", "text", "image_blocks", "warning")

    def __init__(
        self,
        path: Path,
        category: str,
        text: str | None = None,
        image_blocks: list[dict] | None = None,
        warning: str | None = None,
    ):
        self.path = path
        self.category = category
        self.text = text
        self.image_blocks = image_blocks
        self.warning = warning


def scan(input_path: str | Path) -> list[InputFile]:
    """Scan a file or directory and extract content from all recognized files.

    Returns a list of InputFile objects with extracted text or image data.
    Unrecognized files produce a warning but are not fatal.
    """
    path = Path(input_path).resolve()
    if not path.exists():
        raise FileNotFoundError(f"Input path does not exist: {path}")

    if path.is_file():
        files = [path]
    elif path.is_dir():
        files = sorted(
            f for f in path.rglob("*") if f.is_file() and not f.name.startswith(".")
        )
    else:
        raise ValueError(f"Input path is neither a file nor a directory: {path}")

    if not files:
        raise ValueError(f"No files found at: {path}")

    results: list[InputFile] = []
    for fp in files:
        category = _classify(fp)

        if category == "text":
            results.append(InputFile(fp, category, text=_extract_text(fp)))
        elif category == "pdf":
            text, images = _extract_pdf(fp)
            results.append(InputFile(fp, category, text=text, image_blocks=images))
        elif category == "docx":
            results.append(InputFile(fp, category, text=_extract_docx(fp)))
        elif category == "pptx":
            results.append(InputFile(fp, category, text=_extract_pptx(fp)))
        elif category == "image":
            results.append(
                InputFile(fp, category, image_blocks=[_load_image_block(fp)])
            )
        else:
            ext = fp.suffix.lower()
            hint = _UNSUPPORTED_KNOWN.get(ext, "Unknown file type — skipping")
            results.append(InputFile(fp, "unsupported", warning=hint))

    return results


def build_extracted_text(
    input_files: list[InputFile],
    extra_prompt: str | None = None,
) -> str:
    """Build the extracted text document consumed by the intake step."""
    text_parts: list[str] = []
    warnings: list[str] = []

    for f in input_files:
        if f.warning:
            warnings.append(f"- {f.path.name}: {f.warning}")
            continue
        if f.text:
            text_parts.append(f"### File: {f.path.name}\n\n{f.text}")
        if f.image_blocks:
            n_pages = len(f.image_blocks)
            label = f"{n_pages} page image(s)" if n_pages > 1 else "image"
            text_parts.append(
                f"### Image: {f.path.name}\n\n"
                f"(Scanned document with {label}; text extraction was sparse.)"
            )

    extracted_text = "\n\n".join(text_parts) if text_parts else "(no text extracted)"
    if warnings:
        extracted_text += "\n\n## Warnings\n\n" + "\n".join(warnings)
    if extra_prompt:
        extracted_text += f"\n\n## Additional Instructions\n\n{extra_prompt}"

    return extracted_text
